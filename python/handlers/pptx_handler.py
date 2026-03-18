"""PPTX handler using python-pptx."""

from pptx import Presentation
from pptx.util import Inches


def analyze_summary(file_path):
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        shapes_info = []
        for j, shape in enumerate(slide.shapes):
            info = {
                "index": j,
                "shape_type": str(shape.shape_type),
                "name": shape.name,
                "has_text": shape.has_text_frame,
            }
            if shape.has_text_frame:
                info["text"] = shape.text_frame.text[:200] + (
                    "..." if len(shape.text_frame.text) > 200 else ""
                )
            shapes_info.append(info)

        slide_info = {
            "index": i,
            "slide_number": i + 1,
            "shape_count": len(slide.shapes),
            "shapes": shapes_info,
            "has_notes": bool(
                slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip()
            ),
        }
        slides.append(slide_info)

    return {
        "doc_type": "pptx",
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": slides,
    }


def analyze_search(file_path, query):
    prs = Presentation(file_path)
    query_lower = query.lower()
    results = []

    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                if query_lower in shape.text_frame.text.lower():
                    results.append({
                        "type": "shape_text",
                        "slide_index": si,
                        "shape_index": shi,
                        "shape_name": shape.name,
                        "text": shape.text_frame.text[:300],
                    })

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if query_lower in notes_text.lower():
                results.append({
                    "type": "notes",
                    "slide_index": si,
                    "text": notes_text[:300],
                })

    return {"query": query, "matches": results}


def apply_operations(file_path, operations):
    prs = Presentation(file_path)
    applied = []

    for idx, op in enumerate(operations):
        op_type = op["op"]
        try:
            if op_type == "replace_text":
                _op_replace_text(prs, op)
            elif op_type == "add_text_slide":
                _op_add_text_slide(prs, op)
            elif op_type == "delete_slide":
                _op_delete_slide(prs, op)
            else:
                raise ValueError(f"Unknown operation: {op_type}")
            applied.append({"index": idx, "op": op_type, "status": "ok"})
        except Exception as e:
            return {
                "success": False,
                "failed_op_index": idx,
                "failed_op": op_type,
                "error": str(e),
                "applied_before_failure": applied,
            }

    prs.save(file_path)
    return {"success": True, "applied": applied}


def _op_replace_text(prs, op):
    """Locator-based replace: requires target.slide_index + target.shape_index."""
    target = op.get("target")
    if not target or "slide_index" not in target or "shape_index" not in target:
        raise ValueError("replace_text requires target.slide_index and target.shape_index")

    slide_index = target["slide_index"]
    shape_index = target["shape_index"]
    old_text = op["args"]["old_text"]
    new_text = op["args"]["new_text"]

    slides = list(prs.slides)
    if slide_index < 0 or slide_index >= len(slides):
        raise ValueError(f"Invalid slide index: {slide_index} (total: {len(slides)})")

    slide = slides[slide_index]
    shapes = list(slide.shapes)
    if shape_index < 0 or shape_index >= len(shapes):
        raise ValueError(f"Invalid shape index: {shape_index} (total: {len(shapes)} on slide {slide_index})")

    shape = shapes[shape_index]
    if not shape.has_text_frame:
        raise ValueError(f"Shape {shape_index} on slide {slide_index} has no text frame")

    found = False
    for paragraph in shape.text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)
        if old_text in full_text:
            _replace_in_runs(paragraph.runs, old_text, new_text)
            found = True

    if not found:
        raise ValueError(
            f"Text not found in slide {slide_index}, shape {shape_index}: {old_text[:50]}"
        )


def _replace_in_runs(runs, old_text, new_text):
    """Replace text across runs while preserving formatting."""
    combined = ""
    run_map = []
    for i, run in enumerate(runs):
        start = len(combined)
        combined += run.text
        run_map.append((i, start, len(combined)))

    idx = combined.find(old_text)
    if idx == -1:
        return

    end_idx = idx + len(old_text)

    for run_i, char_start, char_end in run_map:
        if char_end <= idx or char_start >= end_idx:
            continue
        run = runs[run_i]
        run_text = run.text
        overlap_start = max(idx, char_start) - char_start
        overlap_end = min(end_idx, char_end) - char_start
        new_run_text = run_text[:overlap_start] + run_text[overlap_end:]
        if char_start <= idx < char_end:
            new_run_text = run_text[:overlap_start] + new_text + run_text[overlap_end:]
        run.text = new_run_text


def _op_add_text_slide(prs, op):
    title = op["args"]["title"]
    body = op["args"].get("body", "")
    insert_at = op.get("target", {}).get("slide_index")

    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = title

    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            shape.text = body
            break

    if insert_at is not None:
        slide_count = len(prs.slides)
        if insert_at < 0 or insert_at > slide_count - 1:
            raise ValueError(f"Invalid slide_index: {insert_at}")
        xml_slides = prs.slides._sldIdLst
        slides_list = list(xml_slides)
        new_slide_elem = slides_list[-1]
        xml_slides.remove(new_slide_elem)
        xml_slides.insert(insert_at, new_slide_elem)


def _op_delete_slide(prs, op):
    target_index = op["target"]["slide_index"]
    slides = prs.slides
    slide_count = len(slides)
    if target_index < 0 or target_index >= slide_count:
        raise ValueError(f"Invalid slide index: {target_index} (total: {slide_count})")

    rId = slides._sldIdLst[target_index].get("r:id") or slides._sldIdLst[target_index].attrib.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rId)
    slides._sldIdLst.remove(slides._sldIdLst[target_index])
