#!/usr/bin/env python3
"""Create minimal DOCX and PPTX test fixtures."""

import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches

FIXTURES_DIR = os.path.join(os.path.dirname(__file__), "fixtures")

# -- DOCX fixture --
doc = Document()
doc.add_heading("Test Document", level=0)
doc.add_paragraph("This is the first paragraph with some sample text.")
doc.add_paragraph("This is the second paragraph for testing replacement.")
doc.add_paragraph("This paragraph contains the keyword searchable content.")

table = doc.add_table(rows=2, cols=2)
table.cell(0, 0).text = "Header A"
table.cell(0, 1).text = "Header B"
table.cell(1, 0).text = "Cell A1"
table.cell(1, 1).text = "Cell B1"

doc.save(os.path.join(FIXTURES_DIR, "test.docx"))
print("Created test.docx")

# -- PPTX fixture --
prs = Presentation()
layout = prs.slide_layouts[1]  # Title and Content

slide1 = prs.slides.add_slide(layout)
slide1.shapes.title.text = "Slide One Title"
for shape in slide1.placeholders:
    if shape.placeholder_format.idx != 0:
        shape.text = "This is the body text of slide one."
        break

slide2 = prs.slides.add_slide(layout)
slide2.shapes.title.text = "Slide Two Title"
for shape in slide2.placeholders:
    if shape.placeholder_format.idx != 0:
        shape.text = "Second slide with searchable keyword here."
        break

prs.save(os.path.join(FIXTURES_DIR, "test.pptx"))
print("Created test.pptx")
